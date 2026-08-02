"""
pptx_builder.py
Сборка готового .pptx файла из сгенерированного контента слайдов,
с фоновыми изображениями (или градиентной заглушкой) и оформлением по стилю.
"""

import io
import logging
from typing import List, Optional

import httpx
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger("slideforge.pptx_builder")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Пресеты стилей: (фон, цвет заголовка, цвет текста, шрифт)
STYLE_PRESETS = {
    "minimal": {"bg": "FFFFFF", "title": "111111", "text": "333333", "accent": "6366F1", "font": "Helvetica"},
    "academic": {"bg": "F7F5F0", "title": "1F2937", "text": "374151", "accent": "1D4ED8", "font": "Georgia"},
    "creative": {"bg": "FDF2F8", "title": "831843", "text": "4B1D3F", "accent": "EC4899", "font": "Verdana"},
    "corporate": {"bg": "FFFFFF", "title": "0F172A", "text": "1E293B", "accent": "0EA5E9", "font": "Calibri"},
    "dark": {"bg": "0F0F14", "title": "FFFFFF", "text": "D1D5DB", "accent": "8B5CF6", "font": "Helvetica"},
}


def _hex_to_rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


async def _download_image(url: str) -> Optional[bytes]:
    try:
        async with httpx.AsyncClient(timeout=20.0) as client:
            resp = await client.get(url)
            if resp.status_code == 200:
                return resp.content
    except httpx.RequestError as exc:
        logger.warning("Не удалось скачать изображение %s: %s", url, exc)
    return None


def _add_gradient_background(slide, color1: str, color2: str):
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    stops = fill.gradient_stops
    stops[0].color.rgb = _hex_to_rgb(color1.lstrip("#"))
    stops[1].color.rgb = _hex_to_rgb(color2.lstrip("#"))
    fill.gradient_angle = 45.0


def _add_solid_background(slide, color_hex: str):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(color_hex)


async def build_pptx(
    title: str,
    slides: List[dict],
    style: str = "minimal",
    language: str = "ru",
) -> bytes:
    """
    slides: список dict вида
        {"title": str, "bullets": [str], "image_url": Optional[str], "notes": Optional[str]}
    Возвращает bytes готового .pptx файла.
    """
    preset = STYLE_PRESETS.get(style, STYLE_PRESETS["minimal"])
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]

    for idx, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)
        is_title_slide = idx == 0
        image_bytes = None

        if slide_data.get("image_url"):
            image_bytes = await _download_image(slide_data["image_url"])

        if image_bytes:
            slide.shapes.add_picture(
                io.BytesIO(image_bytes), 0, 0, width=SLIDE_W, height=SLIDE_H
            )
            # полупрозрачная плашка под текст для читаемости
            overlay = slide.shapes.add_shape(1, 0, Inches(4.6), SLIDE_W, Inches(2.9))
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = _hex_to_rgb("000000")
            overlay.fill.transparency = 0.35
            overlay.line.fill.background()
            title_color = "FFFFFF"
            text_color = "F3F4F6"
            title_top = Inches(4.8)
        else:
            gradient_color2 = preset["accent"]
            _add_gradient_background(slide, preset["bg"], gradient_color2) if is_title_slide else _add_solid_background(
                slide, preset["bg"]
            )
            title_color = preset["title"] if not is_title_slide else "FFFFFF"
            text_color = preset["text"] if not is_title_slide else "F3F4F6"
            title_top = Inches(0.6) if not is_title_slide else Inches(3.0)

        # Заголовок
        title_box = slide.shapes.add_textbox(Inches(0.7), title_top, Inches(11.9), Inches(1.2))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data["title"]
        p.font.size = Pt(40) if is_title_slide else Pt(32)
        p.font.bold = True
        p.font.name = preset["font"]
        p.font.color.rgb = _hex_to_rgb(title_color)
        if is_title_slide:
            p.alignment = PP_ALIGN.CENTER

        # Буллеты
        bullets = slide_data.get("bullets") or []
        if bullets and not is_title_slide:
            body_top = title_top + Inches(1.3)
            body_box = slide.shapes.add_textbox(Inches(0.7), body_top, Inches(11.9), Inches(2.2))
            body_tf = body_box.text_frame
            body_tf.word_wrap = True
            for i, bullet in enumerate(bullets):
                bp = body_tf.paragraphs[0] if i == 0 else body_tf.add_paragraph()
                bp.text = f"•  {bullet}"
                bp.font.size = Pt(20)
                bp.font.name = preset["font"]
                bp.font.color.rgb = _hex_to_rgb(text_color)
                bp.space_after = Pt(10)
        elif bullets and is_title_slide:
            sub_box = slide.shapes.add_textbox(Inches(1.5), title_top + Inches(1.3), Inches(10.3), Inches(1.0))
            sub_tf = sub_box.text_frame
            sub_tf.word_wrap = True
            sp = sub_tf.paragraphs[0]
            sp.text = bullets[0]
            sp.font.size = Pt(20)
            sp.font.name = preset["font"]
            sp.font.color.rgb = _hex_to_rgb(text_color)
            sp.alignment = PP_ALIGN.CENTER

        # Заметки докладчика
        if slide_data.get("notes"):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data["notes"]

        # Номер слайда
        if not is_title_slide:
            num_box = slide.shapes.add_textbox(SLIDE_W - Inches(1.0), SLIDE_H - Inches(0.5), Inches(0.8), Inches(0.4))
            np_ = num_box.text_frame.paragraphs[0]
            np_.text = str(idx + 1)
            np_.font.size = Pt(12)
            np_.font.color.rgb = _hex_to_rgb(text_color)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
